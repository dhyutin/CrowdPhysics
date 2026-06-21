#!/usr/bin/env python3
"""
Build the CrowdPhysics demo deck → devpost/CrowdPhysics_Demo.pptx

Dark, brand-matched 16:9 slides that reuse the architecture diagrams already in
this folder. Re-run after editing:  python devpost/make_deck.py
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
OUT = HERE / "CrowdPhysics_Demo.pptx"
LOGO = ROOT / "brand" / "crowd_physics_logo.png"

IMG_MONITOR = HERE / "monitoring_pipeline_architecture.png"
IMG_DECISION = HERE / "monitor_decision_framework.png"
IMG_SIM = HERE / "simulation_pipeline_architecture.png"
IMG_BRIDGE = HERE / "sim_to_raft_bridge.png"

# ── palette ───────────────────────────────────────────────────────────────────
# Match the logo's own background so it blends seamlessly.
_logo_bg = Image.open(LOGO).convert("RGB").getpixel((6, 6))
BG = RGBColor(*_logo_bg)
CARD = RGBColor(0x16, 0x1C, 0x2A)
CARD2 = RGBColor(0x1B, 0x22, 0x33)
WHITE = RGBColor(0xE9, 0xEE, 0xF5)
MUTE = RGBColor(0x93, 0xA1, 0xB3)
LAV = RGBColor(0xCB, 0xA6, 0xF0)       # light lavender accent
LAV_D = RGBColor(0x8A, 0x63, 0xC4)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
EMER = RGBColor(0x4F, 0xC9, 0x6A)
AMBER = RGBColor(0xE0, 0xA7, 0x32)
CRIM = RGBColor(0xF8, 0x51, 0x49)
BLUE = RGBColor(0x59, 0x9B, 0xF0)

FONT = "Arial"
FONT_H = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


# ── helpers ───────────────────────────────────────────────────────────────────
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, *, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE,
         radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(s, x, y, w, h, runs, *, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.12, italic=False):
    """runs: a string, or list of paragraphs. Each paragraph is a string or a
    list of (text, opts) tuples for inline styling."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        segs = para if isinstance(para, list) else [(para, {})]
        for seg_text, opts in segs:
            r = p.add_run()
            r.text = seg_text
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", bold)
            r.font.italic = opts.get("italic", italic)
            r.font.name = opts.get("font", font)
            r.font.color.rgb = opts.get("color", color)
        if "space_after" in (segs[0][1] if isinstance(segs[0], tuple) else {}):
            pass
    return tb


def kicker(s, label, color=LAV):
    rect(s, 0.85, 0.62, 0.16, 0.46, fill=color)
    text(s, 1.12, 0.55, 11.0, 0.6,
         [[(label.upper(), {"size": 14, "bold": True, "color": color, "font": FONT_H})]],
         anchor=MSO_ANCHOR.MIDDLE)


def title(s, t, y=1.12, size=33, color=WHITE):
    text(s, 0.85, y, 11.6, 1.0,
         [[(t, {"size": size, "bold": True, "color": color, "font": FONT_H})]])


def footer(s, idx):
    text(s, 0.85, SH - 0.5, 6.0, 0.35,
         [[("CROWD", {"size": 10, "bold": True, "color": WHITE}),
           ("PHYSICS", {"size": 10, "bold": True, "color": BLUE})]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, SW - 1.6, SH - 0.5, 0.9, 0.35,
         [[(f"{idx:02d}", {"size": 10, "color": MUTE})]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def img_fit(s, path, bx, by, bw, bh, *, border=True):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = bw
    h = w / ar
    if h > bh:
        h = bh
        w = h * ar
    left = bx + (bw - w) / 2
    top = by + (bh - h) / 2
    if border:
        rect(s, left - 0.04, top - 0.04, w + 0.08, h + 0.08, fill=None,
             line=LAV_D, lw=1.2)
    s.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))


def bullets(s, x, y, w, h, items, *, size=18, gap=12, marker="—", mcolor=LAV,
            tcolor=WHITE):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.1
        p.space_after = Pt(gap)
        rm = p.add_run()
        rm.text = f"{marker}  "
        rm.font.size = Pt(size)
        rm.font.bold = True
        rm.font.color.rgb = mcolor
        rm.font.name = FONT
        if isinstance(item, list):
            for seg_text, opts in item:
                r = p.add_run()
                r.text = seg_text
                r.font.size = Pt(opts.get("size", size))
                r.font.bold = opts.get("bold", False)
                r.font.color.rgb = opts.get("color", tcolor)
                r.font.name = FONT
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = tcolor
            r.font.name = FONT
    return tb


# ════════════════════════════════════════════════════════════════════════════
# 1 · TITLE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
iw, ih = Image.open(LOGO).size
lw = 8.2
lh = lw * ih / iw
img_fit(s, str(LOGO), (SW - lw) / 2, 1.05, lw, lh, border=False)
text(s, 0, lh + 1.1, SW, 0.6,
     [[("Live crowd-crush early warning  +  pre-event crowd-flow simulation",
        {"size": 17, "color": MUTE})]], align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# 2 · THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "The problem", CRIM)
title(s, "Crowd management is a hassle")
text(s, 0.85, 2.0, 11.6, 0.8,
     [[("As the numbers grow, the problems compound — and the failures turn deadly.",
        {"size": 19, "color": WHITE})]])
bullets(s, 0.95, 2.95, 7.4, 3.6, [
    [("Kumbh Mela — ", {"bold": True}), ("the largest human gathering on Earth, with repeated fatal crushes", {"color": MUTE})],
    [("Itaewon, Seoul (2022) · Love Parade, Germany (2010) · Astroworld, USA (2021)", {"color": WHITE})],
    [("Concerts, rallies, religious gatherings, transit hubs — ", {"bold": True}), ("every venue, every year", {"color": MUTE})],
    [("All of them were ", {"color": WHITE}), ("planned months in advance", {"bold": True, "color": CRIM}), (" — and still failed", {"color": WHITE})],
], size=17, mcolor=CRIM, gap=14)
# stat callout
rect(s, 8.7, 2.9, 3.75, 2.2, fill=CARD, line=CRIM, lw=1.4,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
text(s, 8.95, 3.05, 3.3, 1.9,
     [[("A crush is a ", {"size": 18, "color": WHITE}),
       ("physics problem", {"size": 18, "bold": True, "color": CRIM})],
      [("before it is a human one.", {"size": 18, "color": WHITE})],
      [("", {"size": 6})],
      [("By the time you see people fall, it's already too late.",
        {"size": 12.5, "color": MUTE, "italic": True})]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
footer(s, 2)


# ════════════════════════════════════════════════════════════════════════════
# 3 · THE IDEA
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "The idea", LAV)
title(s, "What if a camera could read the physics — and the future?")
text(s, 0.85, 2.35, 11.6, 3.0,
     [[("CrowdPhysics turns ", {"size": 22, "color": WHITE}),
       ("any existing camera", {"size": 22, "bold": True, "color": LAV}),
       (" into a crush early-warning system,", {"size": 22, "color": WHITE})],
      [("and lets you ", {"size": 22, "color": WHITE}),
       ("simulate a venue before the event", {"size": 22, "bold": True, "color": TEAL}),
       (" — through the same perception pipeline.", {"size": 22, "color": WHITE})],
      [("", {"size": 14})],
      [("It learns crowd physics ", {"size": 18, "color": MUTE}),
       ("internally", {"size": 18, "bold": True, "color": WHITE}),
       (" — self-supervised, with no disaster data, by learning what \u201cnormal\u201d looks like.",
        {"size": 18, "color": MUTE})]],
     line_spacing=1.25)
footer(s, 3)


# ════════════════════════════════════════════════════════════════════════════
# 4 · HOW IT WORKS — the 4 stages
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "How it works", LAV)
title(s, "Four stages, from pixels to a plain-language call")
stages = [
    ("01", "Optical flow", "RAFT reads the crowd as pure motion — no faces, no tracking", TEAL),
    ("02", "World model", "a self-supervised model learns the physics & forecasts what's next", LAV),
    ("03", "RL decision", "a model-based RL agent picks the intervention, in imagination", BLUE),
    ("04", "Agentic translation", "Claude turns it into a calibrated risk + plain-language action", EMER),
]
cw, ch, gap = 2.78, 3.2, 0.27
x0 = (SW - (cw * 4 + gap * 3)) / 2
for i, (num, t, d, c) in enumerate(stages):
    x = x0 + i * (cw + gap)
    rect(s, x, 2.5, cw, ch, fill=CARD, line=c, lw=1.4,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    text(s, x + 0.22, 2.7, cw - 0.4, 0.6,
         [[(num, {"size": 26, "bold": True, "color": c})]])
    text(s, x + 0.22, 3.4, cw - 0.44, 0.7,
         [[(t, {"size": 16.5, "bold": True, "color": WHITE})]])
    text(s, x + 0.22, 4.1, cw - 0.44, 1.5,
         [[(d, {"size": 12.5, "color": MUTE})]], line_spacing=1.2)
    if i < 3:
        text(s, x + cw - 0.02, 3.7, gap + 0.04, 0.5,
             [[("\u2192", {"size": 20, "bold": True, "color": LAV})]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.85, 6.05, 11.6, 0.6,
     [[("Tools: ", {"size": 13, "bold": True, "color": LAV}),
       ("PyTorch \u00b7 RAFT (torchvision) \u00b7 LSTM world model \u00b7 Dyna + CQL RL \u00b7 Claude (Sonnet)",
        {"size": 13, "color": MUTE})]])
footer(s, 4)


# ════════════════════════════════════════════════════════════════════════════
# 5 · IT LEARNED PHYSICS — probe R²
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Proof", TEAL)
title(s, "It discovered crowd physics on its own")
text(s, 0.85, 1.95, 11.6, 0.6,
     [[("We linearly probed the unlabeled latent space and recovered interpretable physics:",
        {"size": 15.5, "color": MUTE})]])
probes = [
    ("Crowd velocity", 0.83, TEAL),
    ("Turbulence", 0.78, TEAL),
    ("Backward pressure", 0.84, AMBER),
    ("Boundary stress  (the literal mechanism of a crush)", 0.94, EMER),
]
by = 2.75
bar_x, bar_w = 5.2, 5.7
for i, (name, r2, c) in enumerate(probes):
    y = by + i * 0.74
    text(s, 0.9, y - 0.04, 4.2, 0.5,
         [[(name, {"size": 13.5, "color": WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, bar_x, y + 0.04, bar_w, 0.34, fill=CARD2,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    rect(s, bar_x, y + 0.04, bar_w * r2, 0.34, fill=c,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, bar_x + bar_w + 0.12, y - 0.02, 1.0, 0.45,
         [[(f"R\u00b2 {r2:.2f}", {"size": 13, "bold": True, "color": c})]],
         anchor=MSO_ANCHOR.MIDDLE)
rect(s, 0.9, 5.85, 11.5, 0.95, fill=CARD, line=LAV_D, lw=1.2,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
text(s, 1.15, 5.85, 11.0, 0.95,
     [[("Even the latent dimensions we ", {"size": 14, "color": WHITE}),
       ("couldn't", {"size": 14, "italic": True, "color": WHITE}),
       ("  name still separate pre-anomaly frames from calm ones by ", {"size": 14, "color": WHITE}),
       ("1.56\u03c3", {"size": 14, "bold": True, "color": LAV}),
       (" — early-warning signal we don't yet have words for.", {"size": 14, "color": WHITE})]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)


# ════════════════════════════════════════════════════════════════════════════
# 6 · MONITOR MODE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Monitor mode \u00b7 live", CRIM)
title(s, "Warns before the crush forms")
text(s, 0.85, 1.95, 11.6, 0.6,
     [[("Frame \u2192 optical flow \u2192 world model \u2192 \u201csurprise\u201d \u2192 forecast. The crowd is rendered as a CFD pressure field — the people disappear, only the physics remains.",
        {"size": 14, "color": MUTE})]])
img_fit(s, str(IMG_MONITOR), 0.85, 2.6, 11.6, 4.4)
footer(s, 6)


# ════════════════════════════════════════════════════════════════════════════
# 7 · MULTI-AGENT DECISION FRAMEWORK
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Monitor mode \u00b7 the brain", LAV)
title(s, "No single model decides alone")
text(s, 0.85, 1.95, 11.6, 0.6,
     [[("Anomaly status, imagined futures, the statistical trend, the RL intervention and a counterfactual are fused — Claude reasons over all of it.",
        {"size": 14, "color": MUTE})]])
img_fit(s, str(IMG_DECISION), 0.85, 2.55, 11.6, 4.45)
footer(s, 7)


# ════════════════════════════════════════════════════════════════════════════
# 8 · PLAN MODE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Simulate mode \u00b7 pre-event", TEAL)
title(s, "Simulate the crowd before anyone arrives")
text(s, 0.85, 1.95, 11.6, 0.6,
     [[("Upload a photo or video — agents rebuild the venue in 3D, fill it with a simulated crowd, and surface danger zones, Fruin level-of-service and a safe arrangement plan.",
        {"size": 14, "color": MUTE})]])
img_fit(s, str(IMG_SIM), 0.85, 2.6, 11.6, 4.4)
footer(s, 8)


# ════════════════════════════════════════════════════════════════════════════
# 9 · SIM → RAFT BRIDGE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Simulate mode \u00b7 closing the loop", BLUE)
title(s, "Validate a layout through the same eyes that will watch it")
text(s, 0.85, 1.95, 11.6, 0.9,
     [[("We render the simulation as a synthetic-crowd video and run it through the ", {"size": 15, "color": MUTE}),
       ("same RAFT optical-flow extractor used live", {"size": 15, "bold": True, "color": WHITE}),
       (" — previewing the inflow / outflow each entrance & exit should show on the day.",
        {"size": 15, "color": MUTE})]])
img_fit(s, str(IMG_BRIDGE), 0.85, 3.15, 11.6, 3.35)
footer(s, 9)


# ════════════════════════════════════════════════════════════════════════════
# 10 · TECH STACK / TOOLS
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Built with", LAV)
title(s, "The stack")
cards = [
    ("Perception", "PyTorch \u00b7 RAFT (torchvision, fine-tuned)", TEAL),
    ("World model", "LSTM latent dynamics \u00b7 256 \u2192 64-d", LAV),
    ("Decision (RL)", "Dyna + Conservative Q-Learning", BLUE),
    ("Agent / LLM", "Claude (Sonnet) \u00b7 Anthropic", EMER),
    ("Live capture", "Browserbase \u00b7 yt-dlp", AMBER),
    ("Observability", "Arize AX (tracing + evals)", TEAL),
    ("Alerts", "Fetch.ai heartbeat agent", CRIM),
    ("App", "FastAPI \u00b7 Next.js \u00b7 Three.js", LAV),
]
cw, ch = 5.65, 0.98
gx, gy = 0.55, 0.28
x0, y0 = 0.85, 2.25
for i, (t, d, c) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    rect(s, x, y, cw, ch, fill=CARD, line=c, lw=1.3,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    rect(s, x, y, 0.12, ch, fill=c)
    text(s, x + 0.3, y + 0.1, cw - 0.5, 0.4,
         [[(t, {"size": 14.5, "bold": True, "color": WHITE})]])
    text(s, x + 0.3, y + 0.5, cw - 0.5, 0.4,
         [[(d, {"size": 12, "color": MUTE})]])
footer(s, 10)


# ════════════════════════════════════════════════════════════════════════════
# 11 · USE CASES
# ════════════════════════════════════════════════════════════════════════════
s = slide()
kicker(s, "Where it goes next", EMER)
title(s, "One engine, many domains")
text(s, 0.85, 1.9, 11.6, 0.6,
     [[("Anywhere flow becomes pressure becomes risk — retrain only on what \u201cnormal\u201d looks like.",
        {"size": 15, "color": MUTE})]])
uses = [
    ("Religious gatherings", "Kumbh Mela \u00b7 Hajj \u00b7 pilgrimages", LAV),
    ("Stadiums & concerts", "ingress surges \u00b7 mosh pits \u00b7 exits", CRIM),
    ("Transit & metros", "platforms \u00b7 escalators \u00b7 turnstiles", BLUE),
    ("Airports", "security \u00b7 gates \u00b7 baggage halls", TEAL),
    ("Protests & rallies", "open-area density \u00b7 chokepoints", AMBER),
    ("Retail & evacuations", "Black Friday \u00b7 fire egress \u00b7 drills", EMER),
]
cw, ch = 3.7, 1.55
gx, gy = 0.28, 0.28
x0, y0 = 0.85, 2.7
for i, (t, d, c) in enumerate(uses):
    col = i % 3
    row = i // 3
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    rect(s, x, y, cw, ch, fill=CARD, line=c, lw=1.3,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    text(s, x + 0.28, y + 0.22, cw - 0.5, 0.5,
         [[(t, {"size": 15.5, "bold": True, "color": WHITE})]])
    text(s, x + 0.28, y + 0.78, cw - 0.5, 0.6,
         [[(d, {"size": 12, "color": MUTE})]])
footer(s, 11)


# ════════════════════════════════════════════════════════════════════════════
# 12 · CLOSING
# ════════════════════════════════════════════════════════════════════════════
s = slide()
iw, ih = Image.open(LOGO).size
lw = 5.6
lh = lw * ih / iw
img_fit(s, str(LOGO), (SW - lw) / 2, 1.25, lw, lh, border=False)
text(s, 0, 1.25 + lh + 0.3, SW, 0.9,
     [[("Plan safe.  Monitor live.  ", {"size": 24, "bold": True, "color": WHITE}),
       ("Never react.", {"size": 24, "bold": True, "color": LAV})]],
     align=PP_ALIGN.CENTER)
text(s, 0, 1.25 + lh + 1.15, SW, 0.6,
     [[("Detect danger you never trained on — by learning what normal looks like.",
        {"size": 15, "color": MUTE, "italic": True})]],
     align=PP_ALIGN.CENTER)


prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
